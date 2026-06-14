"""
PgAppForge Investor Pitch Deck Generator
Generates a professional 18-slide PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.util as util
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ─── Color Palette ────────────────────────────────────────────────────────────
PRIMARY_BLUE   = RGBColor(26,  86,  219)
DARK           = RGBColor(15,  23,  42)
ACCENT_PURPLE  = RGBColor(124, 58,  237)
SUCCESS_GREEN  = RGBColor(5,   150, 105)
WARNING_ORANGE = RGBColor(234, 88,  12)
LIGHT_BG       = RGBColor(248, 250, 252)
WHITE          = RGBColor(255, 255, 255)
LIGHT_BLUE     = RGBColor(219, 234, 254)
MID_BLUE       = RGBColor(59,  130, 246)
SLATE          = RGBColor(100, 116, 139)
DARK_SLATE     = RGBColor(51,  65,  85)

# ─── Slide Dimensions (16:9) ──────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


# ─── Shape / Text helpers ─────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT,
                 wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_paragraph(tf, text, font_size=18, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, space_before=0, font_name="Calibri",
                  italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic
    return p


def set_first_paragraph(tf, text, font_size=18, bold=False, color=WHITE,
                        align=PP_ALIGN.LEFT, font_name="Calibri", italic=False):
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic
    return p


def slide_header(slide, title_text, subtitle=None, bg_color=DARK,
                 title_color=WHITE, accent_bar=True):
    """Add a standard header band with title."""
    # Header band
    add_rect(slide, 0, 0, 13.33, 1.4, fill_color=bg_color)
    if accent_bar:
        add_rect(slide, 0, 1.35, 13.33, 0.06, fill_color=PRIMARY_BLUE)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.33), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    set_first_paragraph(tf, title_text, font_size=32, bold=True,
                        color=title_color, font_name="Calibri")
    if subtitle:
        add_paragraph(tf, subtitle, font_size=16, color=LIGHT_BLUE, font_name="Calibri")


def light_bg_slide(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)


def dark_bg_slide(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK)


# ─── Individual Slide Builders ────────────────────────────────────────────────

def slide_01_cover(slide):
    """Cover — dark gradient blue background."""
    # Background layers for depth
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=RGBColor(10, 40, 100))

    # Decorative accent bar top
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=PRIMARY_BLUE)

    # Decorative angled stripe (simulated with a wide tall rect offset)
    add_rect(slide, 8.5, 0, 5.0, 7.5, fill_color=RGBColor(20, 60, 140))

    # Vertical accent left
    add_rect(slide, 0, 0, 0.12, 7.5, fill_color=ACCENT_PURPLE)

    # Seed tag pill
    tag_box = add_rect(slide, 0.6, 1.1, 2.2, 0.42,
                       fill_color=ACCENT_PURPLE)
    add_text_box(slide, "  SEED ROUND 2026", 0.6, 1.12, 2.2, 0.38,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Main title
    add_text_box(slide, "PgAppForge", 0.6, 1.75, 9.0, 1.8,
                 font_size=72, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri")

    # Subtitle
    add_text_box(slide,
                 "The Composable ERP Platform\nfor African Markets",
                 0.6, 3.5, 9.0, 1.6,
                 font_size=30, bold=False, color=LIGHT_BLUE,
                 align=PP_ALIGN.LEFT, font_name="Calibri")

    # Tagline
    add_text_box(slide, "Compose Everything. Deploy Anywhere. Africa First.",
                 0.6, 5.3, 9.0, 0.7,
                 font_size=17, italic=True, color=RGBColor(148, 163, 184),
                 align=PP_ALIGN.LEFT, font_name="Calibri")

    # Bottom bar
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_color=RGBColor(20, 60, 140))
    add_text_box(slide, "Confidential — For Qualified Investors Only",
                 0.5, 7.17, 12.0, 0.28,
                 font_size=10, color=RGBColor(148, 163, 184),
                 align=PP_ALIGN.CENTER)


def slide_02_opportunity(slide):
    """The Opportunity — 3 large stats."""
    dark_bg_slide(slide)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=PRIMARY_BLUE)

    # Title
    add_text_box(slide, "$2.3B Market. 80% Underserved.",
                 0.5, 0.25, 12.3, 0.85,
                 font_size=36, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Divider
    add_rect(slide, 0.5, 1.15, 12.33, 0.04, fill_color=PRIMARY_BLUE)

    stats = [
        ("$2.3B",   "African ERP market,\ngrowing 18% annually",   PRIMARY_BLUE),
        ("80%",     "of African SMEs use\nspreadsheets",           WARNING_ORANGE),
        ("$1T+",    "annual mobile money\ntransaction volume,\nnone integrated with ERP", SUCCESS_GREEN),
    ]

    card_w = 3.8
    gap = 0.37
    start_x = 0.55

    for i, (num, label, color) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        # Card bg
        add_rect(slide, x, 1.4, card_w, 5.5,
                 fill_color=RGBColor(20, 30, 60))
        # Top accent
        add_rect(slide, x, 1.4, card_w, 0.1, fill_color=color)

        # Big number
        add_text_box(slide, num,
                     x, 1.9, card_w, 2.0,
                     font_size=72, bold=True, color=color,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        # Label
        add_text_box(slide, label,
                     x + 0.2, 3.9, card_w - 0.4, 2.6,
                     font_size=20, color=LIGHT_BLUE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")


def slide_03_compliance(slide):
    """Compliance Tsunami — timeline."""
    dark_bg_slide(slide)
    slide_header(slide, "African Governments Mandate Real-Time Compliance",
                 bg_color=DARK)

    # Timeline items
    events = [
        ("2022", "Kenya",    "KRA eTIMS",       "Real-time invoice submission", PRIMARY_BLUE),
        ("2023", "Uganda",   "URA EFRIS",       "Electronic fiscal receipts",   ACCENT_PURPLE),
        ("2024", "Zambia",   "ZRA Smart Invoice","Intelligent invoice system",   SUCCESS_GREEN),
        ("2025+","Nigeria, Tanzania,\nRwanda",
                 "FIRS / TRA / RRA",             "Following suit fast",          WARNING_ORANGE),
    ]

    line_y = 4.2
    add_rect(slide, 0.8, line_y, 11.8, 0.05, fill_color=PRIMARY_BLUE)

    node_x_positions = [1.3, 4.1, 6.9, 9.7]

    for i, (year, country, mandate, desc, color) in enumerate(events):
        x = node_x_positions[i]
        # Connector dot
        add_rect(slide, x + 0.1, line_y - 0.14, 0.28, 0.28, fill_color=color)

        # Year label above line
        add_text_box(slide, year, x - 0.1, line_y - 0.75, 1.0, 0.45,
                     font_size=18, bold=True, color=color,
                     align=PP_ALIGN.CENTER)

        # Card below line
        add_rect(slide, x - 0.3, line_y + 0.35, 2.5, 2.45,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, x - 0.3, line_y + 0.35, 2.5, 0.07, fill_color=color)

        add_text_box(slide, country, x - 0.2, line_y + 0.5, 2.3, 0.5,
                     font_size=15, bold=True, color=WHITE)
        add_text_box(slide, mandate, x - 0.2, line_y + 0.95, 2.3, 0.55,
                     font_size=13, bold=False, color=color)
        add_text_box(slide, desc, x - 0.2, line_y + 1.5, 2.3, 0.9,
                     font_size=11, color=RGBColor(148, 163, 184))

    # Above the line: also add the top context cards
    ctx = [
        ("2022", "Kenya KRA eTIMS", PRIMARY_BLUE),
        ("2023", "Uganda URA EFRIS", ACCENT_PURPLE),
        ("2024", "Zambia ZRA Smart Invoice", SUCCESS_GREEN),
        ("2025+", "NG/TZ/RW mandates", WARNING_ORANGE),
    ]

    # Closing line
    add_rect(slide, 0.4, 6.95, 12.5, 0.04, fill_color=ACCENT_PURPLE)
    add_text_box(slide,
                 "Every business needs ERP that speaks their tax authority's API.",
                 0.5, 6.98, 12.0, 0.45,
                 font_size=16, bold=True, italic=True,
                 color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def slide_04_incumbent_gap(slide):
    """Why Existing Solutions Fail."""
    light_bg_slide(slide)
    slide_header(slide, "The Incumbent Gap",
                 subtitle="Why existing solutions leave Africa behind",
                 bg_color=DARK)

    boxes = [
        ("SAP / Oracle",
         "• $50K+ implementation cost\n• No Africa fintech connectors\n• 18-month deployment timelines\n• Requires specialist consultants",
         WARNING_ORANGE),
        ("Odoo",
         "• Good module ecosystem\n• Sync-only, no event streaming\n• No AI-native capabilities\n• Limited Africa fintech depth",
         ACCENT_PURPLE),
        ("ERPNext",
         "• Open source but shallow\n• No composable policy algebra\n• No composable AI pipelines\n• Hard to extend without forks",
         PRIMARY_BLUE),
        ("Spreadsheets",
         "• Where 80% of SMEs still live\n• Zero compliance capabilities\n• No scale, no audit trail\n• Manual reconciliation nightmare",
         SUCCESS_GREEN),
    ]

    bw = 2.9
    bh = 4.7
    gap = 0.3
    start_x = 0.4

    for i, (title, body, color) in enumerate(boxes):
        x = start_x + i * (bw + gap)
        add_rect(slide, x, 1.55, bw, bh,
                 fill_color=WHITE,
                 line_color=RGBColor(226, 232, 240), line_width=0.5)
        add_rect(slide, x, 1.55, bw, 0.12, fill_color=color)
        add_text_box(slide, title,
                     x + 0.15, 1.7, bw - 0.3, 0.65,
                     font_size=18, bold=True, color=DARK,
                     font_name="Calibri")
        add_text_box(slide, body,
                     x + 0.15, 2.4, bw - 0.3, 3.5,
                     font_size=14, color=DARK_SLATE,
                     font_name="Calibri")


def slide_05_introducing(slide):
    """Introducing PgAppForge."""
    dark_bg_slide(slide)
    slide_header(slide, "Composable ERP Built for Africa's Stack",
                 subtitle="PgAppForge — the platform that speaks every African API",
                 bg_color=DARK)

    pillars = [
        ("📦", "155 Domain Modules",
         "Finance, GRC, HCM, CRM,\nOperations, Platform — all\ncomposable, all production-ready",
         PRIMARY_BLUE),
        ("🌍", "Africa Connectors",
         "MTN MoMo, Airtel Money,\nM-Pesa, eTIMS, EFRIS, ZRA\nout of the box",
         SUCCESS_GREEN),
        ("🔗", "Composability",
         "11 primitives — compose\nany vertical without coupling\nor forking upstream",
         ACCENT_PURPLE),
        ("🤖", "AI-Native",
         "LiteLLM, RAG, NL analytics,\ncomposable agent pipelines\nbuilt into the core",
         WARNING_ORANGE),
    ]

    pw = 2.9
    ph = 4.8
    gap = 0.3
    sx = 0.4

    for i, (icon, title, body, color) in enumerate(pillars):
        x = sx + i * (pw + gap)
        add_rect(slide, x, 1.55, pw, ph,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, x, 1.55, pw, 0.1, fill_color=color)

        add_text_box(slide, icon,
                     x, 1.75, pw, 0.8,
                     font_size=32, align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     x + 0.1, 2.55, pw - 0.2, 0.7,
                     font_size=17, bold=True, color=color,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        add_text_box(slide, body,
                     x + 0.15, 3.3, pw - 0.3, 2.8,
                     font_size=13, color=LIGHT_BLUE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")


def slide_06_composability_moat(slide):
    """11 Primitives Make Everything Possible."""
    dark_bg_slide(slide)
    slide_header(slide, "11 Primitives Make Everything Possible",
                 subtitle="The Composability Moat",
                 bg_color=DARK)

    columns = [
        ("Core Primitives",
         ["Event Router", "Model Mixins", "Policy Algebra", "Sub-workflow"],
         PRIMARY_BLUE),
        ("AI & Data",
         ["AI Pipeline", "Rule→Event Bridge", "Metric Registry", "PDL Extends"],
         ACCENT_PURPLE),
        ("Integration",
         ["View Slots", "GraphQL Federation", "Cross-Tenant"],
         SUCCESS_GREEN),
    ]

    cw = 3.8
    ch = 5.0
    gap = 0.35
    sx = 0.45

    for i, (header, items, color) in enumerate(columns):
        x = sx + i * (cw + gap)
        add_rect(slide, x, 1.55, cw, ch,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, x, 1.55, cw, 0.1, fill_color=color)

        add_text_box(slide, header,
                     x + 0.15, 1.68, cw - 0.3, 0.55,
                     font_size=16, bold=True, color=color,
                     align=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            iy = 2.35 + j * 0.82
            add_rect(slide, x + 0.2, iy, cw - 0.4, 0.62,
                     fill_color=RGBColor(30, 45, 80),
                     line_color=color, line_width=0.5)
            add_text_box(slide, item,
                         x + 0.3, iy + 0.1, cw - 0.6, 0.45,
                         font_size=14, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER)

    # Key claim banner
    add_rect(slide, 0.4, 6.75, 12.5, 0.62, fill_color=RGBColor(26, 86, 219))
    add_text_box(slide,
                 "A developer composes a SACCO + mobile money + tax platform in weeks, not months.",
                 0.6, 6.8, 12.1, 0.52,
                 font_size=15, bold=True, italic=True,
                 color=WHITE, align=PP_ALIGN.CENTER)


def slide_07_tech_moat(slide):
    """Why This Can't Be Replicated Quickly."""
    light_bg_slide(slide)
    slide_header(slide, "Why This Can't Be Replicated Quickly",
                 subtitle="Technology moat across three dimensions",
                 bg_color=DARK)

    sections = [
        ("🔐  Africa Compliance Layer",
         "18–24 months to build from scratch",
         "eTIMS, EFRIS, ZRA, goAML,\nSASRA, CRB integrations",
         WARNING_ORANGE),
        ("📚  588 Capability Models",
         "Domain expertise, not just CRUD",
         "IFRS 16 lease accounting,\nSoD analysis, Kanban lean,\nMES OEE — pre-built",
         ACCENT_PURPLE),
        ("⚙️  Composability Architecture",
         "Deeper than Odoo's _inherit",
         "EventRouter + ModelMixins\n+ PolicyAlgebra + AI Pipeline\n= true composability",
         PRIMARY_BLUE),
    ]

    bw = 3.9
    bh = 5.0
    gap = 0.32
    sx = 0.42

    for i, (title, subtitle_txt, body, color) in enumerate(sections):
        x = sx + i * (bw + gap)
        add_rect(slide, x, 1.55, bw, bh,
                 fill_color=WHITE,
                 line_color=RGBColor(226, 232, 240), line_width=0.5)
        add_rect(slide, x, 1.55, bw, 0.12, fill_color=color)
        add_text_box(slide, title,
                     x + 0.15, 1.72, bw - 0.3, 0.7,
                     font_size=16, bold=True, color=DARK, font_name="Calibri")
        add_text_box(slide, subtitle_txt,
                     x + 0.15, 2.42, bw - 0.3, 0.55,
                     font_size=13, italic=True, color=color, font_name="Calibri")
        add_rect(slide, x + 0.15, 2.97, bw - 0.3, 0.03,
                 fill_color=RGBColor(226, 232, 240))
        add_text_box(slide, body,
                     x + 0.15, 3.05, bw - 0.3, 3.1,
                     font_size=14, color=DARK_SLATE, font_name="Calibri")


def slide_08_sacco(slide):
    """Use Case: SACCO Platform."""
    dark_bg_slide(slide)
    slide_header(slide, "A Complete SACCO in 8 Weeks",
                 subtitle="Use Case — SACCO Platform Composition",
                 bg_color=DARK)

    steps = [
        ("Member\nKYC",          PRIMARY_BLUE),
        ("Loan\nApplication",    ACCENT_PURPLE),
        ("Committee\nApproval",  SUCCESS_GREEN),
        ("MTN MoMo\nDisbursement", WARNING_ORANGE),
        ("eTIMS\nInvoice",       PRIMARY_BLUE),
        ("Repayment\nTracking",  ACCENT_PURPLE),
        ("SASRA\nReporting",     SUCCESS_GREEN),
        ("NL Analytics",        WARNING_ORANGE),
    ]

    sw = 1.45
    sh = 1.3
    sy = 1.65
    gap = 0.12
    sx = 0.25

    for i, (label, color) in enumerate(steps):
        x = sx + i * (sw + gap)
        add_rect(slide, x, sy, sw, sh,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, x, sy, sw, 0.1, fill_color=color)
        add_text_box(slide, label,
                     x + 0.05, sy + 0.2, sw - 0.1, sh - 0.25,
                     font_size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        if i < len(steps) - 1:
            ax = x + sw + 0.01
            add_text_box(slide, "→",
                         ax, sy + 0.35, 0.14, 0.6,
                         font_size=18, bold=True, color=SLATE,
                         align=PP_ALIGN.CENTER)

    # Policy algebra note
    add_text_box(slide, "Policy Algebra",
                 sx + 2 * (sw + gap) - 0.2, sy + sh + 0.15, 1.8, 0.4,
                 font_size=11, color=SUCCESS_GREEN,
                 align=PP_ALIGN.CENTER)
    add_rect(slide, sx + 2 * (sw + gap) + 0.5, sy + sh + 0.08, 0.02, 0.55,
             fill_color=SUCCESS_GREEN)

    # NL Query example
    add_rect(slide, 0.5, 3.45, 12.3, 0.75,
             fill_color=RGBColor(20, 30, 60))
    add_text_box(slide,
                 '💬  NL Query: "What is our NPL ratio?"  →  Instant answer via composable AI analytics',
                 0.65, 3.52, 12.0, 0.58,
                 font_size=15, italic=True, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)

    # Summary banner
    add_rect(slide, 0.5, 4.4, 12.3, 1.2,
             fill_color=RGBColor(20, 50, 100))
    add_rect(slide, 0.5, 4.4, 12.3, 0.07, fill_color=SUCCESS_GREEN)
    add_text_box(slide,
                 "Full SACCO platform. 8 weeks. 2 developers.",
                 0.6, 4.55, 12.0, 0.55,
                 font_size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text_box(slide,
                 "Member registry → lending → disbursement → compliance → analytics — all composed, nothing forked.",
                 0.6, 5.1, 12.0, 0.42,
                 font_size=14, italic=True, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)

    # Arch note
    add_text_box(slide,
                 "Architecture: EventRouter + ModelMixins + PolicyAlgebra + MoMo connector + eTIMS connector",
                 0.5, 5.75, 12.3, 0.45,
                 font_size=12, color=SLATE, align=PP_ALIGN.CENTER)


def slide_09_trade_finance(slide):
    """Use Case: Trade Finance."""
    dark_bg_slide(slide)
    slide_header(slide, "Trade Finance on Top of Core ERP",
                 subtitle="Use Case — Zero-fork vertical extension via ModelMixinRegistry",
                 bg_color=DARK)

    # Left panel: Mixin adds
    add_rect(slide, 0.4, 1.6, 5.8, 5.3,
             fill_color=RGBColor(20, 30, 60))
    add_rect(slide, 0.4, 1.6, 5.8, 0.1, fill_color=ACCENT_PURPLE)

    add_text_box(slide, "ModelMixinRegistry adds to ARInvoice:",
                 0.55, 1.72, 5.5, 0.55,
                 font_size=15, bold=True, color=ACCENT_PURPLE)

    fields = [
        "letter_of_credit_id",
        "lc_expiry_date",
        "presenting_bank",
        "swift_code",
    ]
    for j, f in enumerate(fields):
        fy = 2.32 + j * 0.62
        add_rect(slide, 0.6, fy, 5.4, 0.5,
                 fill_color=RGBColor(30, 45, 90))
        add_text_box(slide, f"  {f}",
                     0.65, fy + 0.07, 5.2, 0.38,
                     font_size=14, bold=True, color=SUCCESS_GREEN,
                     font_name="Courier New")

    add_text_box(slide, "(Finance plugin untouched — no fork required)",
                 0.55, 4.9, 5.5, 0.7,
                 font_size=12, italic=True, color=SLATE)

    # Right panel: Event flow
    add_rect(slide, 6.7, 1.6, 6.2, 5.3,
             fill_color=RGBColor(20, 30, 60))
    add_rect(slide, 6.7, 1.6, 6.2, 0.1, fill_color=PRIMARY_BLUE)

    add_text_box(slide, "Event Flow:",
                 6.85, 1.72, 5.8, 0.55,
                 font_size=15, bold=True, color=PRIMARY_BLUE)

    events = [
        ("shipment_confirmed", "Event emitted by ShipmentView", WARNING_ORANGE),
        ("→", "", WHITE),
        ("LC Drawing Request", "Auto-generated via Rule→Event Bridge", SUCCESS_GREEN),
        ("→", "", WHITE),
        ("Bank Notification", "SWIFT MT700 payload dispatched", ACCENT_PURPLE),
        ("→", "", WHITE),
        ("Finance Records", "ARInvoice updated automatically", PRIMARY_BLUE),
    ]
    ey = 2.3
    for text, sub, color in events:
        if text == "→":
            add_text_box(slide, "↓", 9.3, ey, 0.5, 0.3,
                         font_size=18, bold=True, color=SLATE,
                         align=PP_ALIGN.CENTER)
            ey += 0.28
        else:
            add_rect(slide, 6.85, ey, 5.8, 0.55,
                     fill_color=RGBColor(30, 45, 90))
            add_text_box(slide, text, 7.0, ey + 0.05, 5.5, 0.3,
                         font_size=13, bold=True, color=color)
            if sub:
                add_text_box(slide, sub, 7.0, ey + 0.32, 5.5, 0.28,
                             font_size=10, italic=True, color=SLATE)
            ey += 0.62

    # Bottom claim
    add_rect(slide, 0.4, 7.05, 12.5, 0.35, fill_color=PRIMARY_BLUE)
    add_text_box(slide, "All via composition — the Finance plugin is completely untouched.",
                 0.5, 7.08, 12.2, 0.28,
                 font_size=13, bold=True, italic=True,
                 color=WHITE, align=PP_ALIGN.CENTER)


def slide_10_library(slide):
    """The 588 Capability Library."""
    light_bg_slide(slide)
    slide_header(slide, "588 Pre-Built Capability Models",
                 subtitle="Deep domain coverage across 8 verticals",
                 bg_color=DARK)

    domains = [
        ("Finance",     45, "IFRS 16, hedge accounting,\nmaterial ledger, JV accounting",     PRIMARY_BLUE),
        ("GRC",         28, "SoD analysis, ERM,\nethics hotline, anti-bribery",               WARNING_ORANGE),
        ("HCM",         52, "Recruiting, performance mgmt,\nposition management, payroll",     ACCENT_PURPLE),
        ("CRM",         38, "Loyalty, PRM, territory mgmt,\nCPQ, field service",              SUCCESS_GREEN),
        ("Operations",  67, "Capacity scheduling,\nlean/kanban, MES OEE",                     PRIMARY_BLUE),
        ("Fintech",     89, "Core banking, SACCO,\nmobile money, trade finance",              WARNING_ORANGE),
        ("Platform",    89, "Analytics, process mining,\nEDI, iPaaS, API gateway",            ACCENT_PURPLE),
        ("Industry",    180, "Agri, health, education,\ngovernment, construction, energy",    SUCCESS_GREEN),
    ]

    # Total badge
    add_rect(slide, 10.8, 1.55, 2.1, 1.5,
             fill_color=DARK)
    add_text_box(slide, "588", 10.8, 1.6, 2.1, 0.95,
                 font_size=54, bold=True, color=WARNING_ORANGE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text_box(slide, "Total\nModels", 10.8, 2.52, 2.1, 0.5,
                 font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    row_h = 0.63
    for i, (domain, count, desc, color) in enumerate(domains):
        row = i % 4
        col = i // 4
        x = 0.4 + col * 5.1
        y = 1.55 + row * (row_h + 0.1)

        add_rect(slide, x, y, 4.8, row_h,
                 fill_color=WHITE,
                 line_color=RGBColor(226, 232, 240), line_width=0.5)
        add_rect(slide, x, y, 0.08, row_h, fill_color=color)

        add_text_box(slide, domain,
                     x + 0.2, y + 0.04, 1.2, 0.4,
                     font_size=15, bold=True, color=DARK)
        add_text_box(slide, str(count),
                     x + 1.45, y + 0.04, 0.7, 0.4,
                     font_size=20, bold=True, color=color)
        add_text_box(slide, "models",
                     x + 2.1, y + 0.14, 0.7, 0.28,
                     font_size=10, color=SLATE)
        add_text_box(slide, desc,
                     x + 1.45, y + 0.38, 3.2, 0.3,
                     font_size=9, color=DARK_SLATE)


def slide_11_pdl_studio(slide):
    """Visual PDL Designer."""
    dark_bg_slide(slide)
    slide_header(slide, "The JDL Studio for Africa's ERPs — and More",
                 subtitle="One schema to rule them all",
                 bg_color=DARK)

    # Left: features list
    features = [
        ("Import", "Any of 588 capability models from 8 domains"),
        ("Design", "Drag-and-drop FK relationships, crossing minimization"),
        ("Generate", "YAML → model + migration + REST API + tests + Dockerfile + K8s"),
        ("Deploy",  "Apply migrations in one command, zero manual steps"),
    ]

    for i, (label, desc) in enumerate(features):
        y = 1.7 + i * 1.15
        add_rect(slide, 0.4, y, 6.0, 0.95,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, 0.4, y, 0.1, 0.95, fill_color=ACCENT_PURPLE)
        add_text_box(slide, label, 0.65, y + 0.05, 1.2, 0.45,
                     font_size=16, bold=True, color=ACCENT_PURPLE)
        add_text_box(slide, desc, 0.65, y + 0.48, 5.6, 0.4,
                     font_size=13, color=LIGHT_BLUE)

    # Right: command box
    add_rect(slide, 7.0, 1.6, 5.9, 2.2,
             fill_color=RGBColor(10, 15, 30))
    add_rect(slide, 7.0, 1.6, 5.9, 0.35,
             fill_color=RGBColor(30, 40, 70))
    add_text_box(slide, "  Terminal", 7.05, 1.62, 2.5, 0.3,
                 font_size=11, color=SLATE)
    add_text_box(slide,
                 "$ flask forge gen pdl schema.yaml \\\n"
                 "    --with-k8s \\\n"
                 "    --apply-migrations\n\n"
                 "✓ Models generated (12 files)\n"
                 "✓ REST APIs scaffolded\n"
                 "✓ Tests written\n"
                 "✓ Dockerfile + K8s manifests\n"
                 "✓ Migrations applied",
                 7.15, 2.02, 5.6, 1.68,
                 font_size=11, color=SUCCESS_GREEN,
                 font_name="Courier New")

    # Schema preview
    add_rect(slide, 7.0, 4.0, 5.9, 3.3,
             fill_color=RGBColor(10, 15, 30))
    add_rect(slide, 7.0, 4.0, 5.9, 0.3, fill_color=RGBColor(30, 40, 70))
    add_text_box(slide, "  schema.yaml", 7.05, 4.02, 3.0, 0.26,
                 font_size=11, color=SLATE)
    yaml_text = (
        "entity: LoanApplication\n"
        "extends: [SACCOMixin, KYCMixin]\n"
        "fields:\n"
        "  - name: amount\n"
        "    type: Decimal\n"
        "    validators: [min:1000]\n"
        "  - name: purpose\n"
        "    type: str\n"
        "connectors: [mtn_momo, etims]\n"
        "policy: committee_approval"
    )
    add_text_box(slide, yaml_text,
                 7.1, 4.35, 5.7, 2.85,
                 font_size=10.5, color=LIGHT_BLUE,
                 font_name="Courier New")


def slide_12_business_model(slide):
    """Business Model — Open Core + Commercial Modules."""
    light_bg_slide(slide)
    slide_header(slide, "Open Core + Commercial Modules",
                 subtitle="Multi-stream revenue with viral developer distribution",
                 bg_color=DARK)

    streams = [
        ("Open Source Core",
         "Free forever",
         "Community grows the ecosystem. Developers build on PgAppForge.\nNetwork effects compound.",
         "Community",
         PRIMARY_BLUE),
        ("Compliance Modules",
         "$X / month / tenant",
         "eTIMS, EFRIS, ZRA — required by law.\nHigh retention, recurring, defensible.",
         "Recurring",
         SUCCESS_GREEN),
        ("AI Modules",
         "$X / month",
         "NL analytics, predictive scoring,\ncomposable agent pipelines.",
         "Expansion",
         ACCENT_PURPLE),
        ("Enterprise Support",
         "$X / year",
         "SLA-backed support, dedicated CSM,\npriority security patches.",
         "Enterprise",
         WARNING_ORANGE),
    ]

    bw = 2.9
    bh = 4.7
    gap = 0.3
    sx = 0.4

    for i, (title, price, body, badge, color) in enumerate(streams):
        x = sx + i * (bw + gap)
        add_rect(slide, x, 1.55, bw, bh,
                 fill_color=WHITE,
                 line_color=RGBColor(226, 232, 240), line_width=0.5)
        add_rect(slide, x, 1.55, bw, 0.12, fill_color=color)

        # Badge
        add_rect(slide, x + bw - 1.35, 1.68, 1.25, 0.38,
                 fill_color=color)
        add_text_box(slide, badge, x + bw - 1.3, 1.7, 1.2, 0.32,
                     font_size=10, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     x + 0.15, 1.72, bw - 1.5, 0.55,
                     font_size=15, bold=True, color=DARK)
        add_text_box(slide, price,
                     x + 0.15, 2.28, bw - 0.3, 0.55,
                     font_size=20, bold=True, color=color)
        add_rect(slide, x + 0.15, 2.85, bw - 0.3, 0.03,
                 fill_color=RGBColor(226, 232, 240))
        add_text_box(slide, body,
                     x + 0.15, 2.95, bw - 0.3, 2.9,
                     font_size=13, color=DARK_SLATE)

    # SI Partner footnote
    add_rect(slide, 0.4, 6.5, 12.5, 0.75,
             fill_color=RGBColor(30, 40, 70))
    add_text_box(slide,
                 "SI Partner Program: 20% first-year revenue share — incentivises certified integrators to sell and deploy PgAppForge",
                 0.6, 6.58, 12.0, 0.55,
                 font_size=14, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)


def slide_13_gtm(slide):
    """Go-to-Market."""
    dark_bg_slide(slide)
    slide_header(slide, "Developer-Led to Enterprise",
                 subtitle="Go-to-Market Strategy",
                 bg_color=DARK)

    channels = [
        ("🛠️  Developer-Led",
         "Open source + GitHub\nAfrican dev communities\nHackathons + bootcamps\nDev-first documentation",
         PRIMARY_BLUE),
        ("🤝  SI Partnerships",
         "5 certified integrators\nKenya / Uganda / Tanzania\n/ Zambia coverage\nRevenue share incentives",
         ACCENT_PURPLE),
        ("📱  Vertical SaaS",
         "SACCO platform builders\nTrade finance boutiques\nInsurtech developers\nAgri-fin startups",
         SUCCESS_GREEN),
        ("🏦  Enterprise Direct",
         "Tier-1 banks\nMobile network operators\nGovernment agencies\nInternal tool teams",
         WARNING_ORANGE),
    ]

    cw = 2.9
    ch = 4.8
    gap = 0.3
    sx = 0.4

    for i, (title, body, color) in enumerate(channels):
        x = sx + i * (cw + gap)
        add_rect(slide, x, 1.55, cw, ch,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, x, 1.55, cw, 0.1, fill_color=color)

        add_text_box(slide, title,
                     x + 0.15, 1.7, cw - 0.3, 0.65,
                     font_size=15, bold=True, color=color)

        add_text_box(slide, body,
                     x + 0.15, 2.4, cw - 0.3, 3.55,
                     font_size=13, color=LIGHT_BLUE)

    # Funnel note
    add_rect(slide, 0.4, 6.55, 12.5, 0.75,
             fill_color=RGBColor(20, 40, 80))
    add_text_box(slide,
                 "Developer adoption → SI deployment → Enterprise land-and-expand — each layer feeds the next",
                 0.6, 6.63, 12.0, 0.55,
                 font_size=14, italic=True, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)


def slide_14_traction(slide):
    """Traction — Early Signals."""
    light_bg_slide(slide)
    slide_header(slide, "Early Signals",
                 subtitle="Platform capabilities already built — not vaporware",
                 bg_color=DARK)

    metrics = [
        ("155",    "domain modules across\n8 service families",  "(live)",     PRIMARY_BLUE),
        ("10",     "mobile money &\npayment connectors",          "(live)",     SUCCESS_GREEN),
        ("3",      "African tax authorities\nintegrated",         "KE / UG / ZM", WARNING_ORANGE),
        ("588",    "pre-built capability\nmodels in library",     "(catalogued)", ACCENT_PURPLE),
        ("<10 min","PDL schema → fully\ndeployable service",      "(benchmarked)", PRIMARY_BLUE),
    ]

    mw = 2.25
    mh = 4.5
    gap = 0.22
    sx = 0.4

    for i, (num, label, badge, color) in enumerate(metrics):
        x = sx + i * (mw + gap)
        add_rect(slide, x, 1.55, mw, mh,
                 fill_color=WHITE,
                 line_color=RGBColor(226, 232, 240), line_width=0.5)
        add_rect(slide, x, 1.55, mw, 0.1, fill_color=color)

        fs = 54 if len(num) <= 3 else 36
        add_text_box(slide, num,
                     x, 1.75, mw, 1.3,
                     font_size=fs, bold=True, color=color,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        add_text_box(slide, label,
                     x + 0.1, 3.05, mw - 0.2, 1.35,
                     font_size=13, color=DARK,
                     align=PP_ALIGN.CENTER)

        add_rect(slide, x + 0.2, 4.4, mw - 0.4, 0.38,
                 fill_color=color)
        add_text_box(slide, badge,
                     x + 0.2, 4.43, mw - 0.4, 0.32,
                     font_size=10, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # Footer note
    add_rect(slide, 0.4, 6.35, 12.5, 0.9,
             fill_color=DARK)
    add_text_box(slide,
                 "All capabilities above are implemented and testable today. "
                 "The 588-model library and 10 connectors represent "
                 "genuine engineering investment — not placeholder features.",
                 0.6, 6.42, 12.0, 0.75,
                 font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def slide_15_ask(slide):
    """The Ask — Seed Round."""
    dark_bg_slide(slide)
    slide_header(slide, "Seed Round: $[X]M",
                 subtitle="Use of Funds",
                 bg_color=DARK)

    uses = [
        ("Engineering", "60%",
         "• 5 senior engineers (backend + platform)\n"
         "• Nigeria FIRS + Tanzania TRA compliance modules\n"
         "• AI features: predictive scoring, NL analytics v2\n"
         "• Performance, security hardening",
         PRIMARY_BLUE),
        ("Pilot Deployments", "20%",
         "• 3 anchor customers at cost (KE / UG / ZM)\n"
         "• Customer success embedded support\n"
         "• Case study and reference material development",
         SUCCESS_GREEN),
        ("Partnerships", "10%",
         "• SI certification program (5 partners)\n"
         "• Developer evangelism and community building\n"
         "• Conference presence in Nairobi, Kampala, Lusaka",
         ACCENT_PURPLE),
        ("Operations", "10%",
         "• Legal (IP, entity structure, contracts)\n"
         "• Security audit (SOC 2 readiness)\n"
         "• Cloud infrastructure baseline",
         WARNING_ORANGE),
    ]

    bw = 2.9
    bh = 4.9
    gap = 0.3
    sx = 0.4

    for i, (title, pct, body, color) in enumerate(uses):
        x = sx + i * (bw + gap)
        add_rect(slide, x, 1.55, bw, bh,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, x, 1.55, bw, 0.1, fill_color=color)

        add_text_box(slide, pct,
                     x, 1.7, bw, 0.85,
                     font_size=48, bold=True, color=color,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        add_text_box(slide, title,
                     x + 0.1, 2.6, bw - 0.2, 0.55,
                     font_size=15, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        add_rect(slide, x + 0.2, 3.18, bw - 0.4, 0.03,
                 fill_color=RGBColor(40, 60, 100))
        add_text_box(slide, body,
                     x + 0.15, 3.25, bw - 0.3, 2.9,
                     font_size=12, color=LIGHT_BLUE)


def slide_16_milestones(slide):
    """18-Month Milestones — Path to Series A."""
    light_bg_slide(slide)
    slide_header(slide, "Path to Series A",
                 subtitle="18-Month Milestones",
                 bg_color=DARK)

    milestones = [
        ("Month 6",
         [
             "5 certified SI partners",
             "10 live deployments",
             "Nigeria FIRS compliance module live",
             "Developer community: 500+ GitHub stars",
         ],
         PRIMARY_BLUE),
        ("Month 12",
         [
             "50 live deployments",
             "Pan-East-Africa coverage (KE/UG/TZ/ZM/RW)",
             "$Xk ARR — compliance module revenue",
             "Tanzania TRA compliance live",
         ],
         ACCENT_PURPLE),
        ("Month 18",
         [
             "Series A ready",
             "West Africa expansion (NG/GH)",
             "$XM ARR run rate",
             "20+ certified SI partners",
         ],
         SUCCESS_GREEN),
    ]

    # Timeline bar
    add_rect(slide, 0.8, 4.5, 11.7, 0.08, fill_color=PRIMARY_BLUE)
    nodes = [2.35, 6.15, 10.0]
    labels_x = [1.5, 5.3, 9.2]
    for nx in nodes:
        add_rect(slide, nx - 0.12, 4.4, 0.24, 0.28, fill_color=ACCENT_PURPLE)

    mw = 3.65
    mh = 4.65
    gap = 0.38
    sx = 0.45

    for i, (period, bullets, color) in enumerate(milestones):
        x = sx + i * (mw + gap)
        add_rect(slide, x, 1.55, mw, mh,
                 fill_color=WHITE,
                 line_color=RGBColor(226, 232, 240), line_width=0.5)
        add_rect(slide, x, 1.55, mw, 0.12, fill_color=color)

        add_text_box(slide, period,
                     x + 0.15, 1.7, mw - 0.3, 0.55,
                     font_size=20, bold=True, color=color)

        for j, bullet in enumerate(bullets):
            by = 2.35 + j * 0.78
            add_rect(slide, x + 0.15, by, 0.12, 0.38, fill_color=color)
            add_text_box(slide, bullet,
                         x + 0.38, by, mw - 0.55, 0.55,
                         font_size=13, color=DARK)

    # Series A banner
    add_rect(slide, 0.4, 6.45, 12.5, 0.8,
             fill_color=DARK)
    add_rect(slide, 0.4, 6.45, 12.5, 0.07, fill_color=SUCCESS_GREEN)
    add_text_box(slide,
                 "Target: Series A metrics achieved by Month 18 — $XM ARR, 50+ live deployments, West Africa beachhead",
                 0.6, 6.56, 12.0, 0.55,
                 font_size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)


def slide_17_why_now(slide):
    """Why Now — The Compliance Window Is Open."""
    dark_bg_slide(slide)
    slide_header(slide, "The Compliance Window Is Open",
                 subtitle="Why Now — three converging forces",
                 bg_color=DARK)

    trends = [
        ("Tax Mandate\nExpansion",
         "Real-time compliance required by law — businesses MUST act.\n"
         "Every new mandate is a forcing function to upgrade from spreadsheets.\n"
         "KE, UG, ZM live. NG, TZ, RW imminent.",
         "5 countries\nmandated",
         WARNING_ORANGE),
        ("Mobile Money\nMaturity",
         "MTN MoMo, Airtel Money, M-Pesa APIs are stable,\nwell-documented, and at scale.\n"
         "$1T+ annual transaction volume with no ERP integration layer.",
         "$1T+\nannually",
         SUCCESS_GREEN),
        ("LLM Cost\nCollapse",
         "AI-powered ERP analytics are now affordable for African SMEs.\n"
         "$0.001/query means NL business intelligence is viable.\n"
         "PgAppForge ships AI-native from day one.",
         "$0.001\nper query",
         ACCENT_PURPLE),
    ]

    tw = 3.8
    th = 4.9
    gap = 0.33
    sx = 0.42

    for i, (title, body, stat, color) in enumerate(trends):
        x = sx + i * (tw + gap)
        add_rect(slide, x, 1.55, tw, th,
                 fill_color=RGBColor(20, 30, 60))
        add_rect(slide, x, 1.55, tw, 0.1, fill_color=color)

        add_text_box(slide, title,
                     x + 0.15, 1.7, tw - 0.3, 0.7,
                     font_size=18, bold=True, color=color)

        add_text_box(slide, body,
                     x + 0.15, 2.45, tw - 0.3, 2.2,
                     font_size=13, color=LIGHT_BLUE)

        add_rect(slide, x + 0.15, 4.65, tw - 0.3, 0.65,
                 fill_color=RGBColor(30, 45, 90))
        add_text_box(slide, stat,
                     x + 0.15, 4.68, tw - 0.3, 0.58,
                     font_size=16, bold=True, color=color,
                     align=PP_ALIGN.CENTER)

    # Urgency banner
    add_rect(slide, 0.4, 6.65, 12.5, 0.68,
             fill_color=RGBColor(26, 86, 219))
    add_text_box(slide,
                 "18–24 months to establish market position before SAP / Oracle respond.",
                 0.6, 6.72, 12.0, 0.52,
                 font_size=16, bold=True, italic=True,
                 color=WHITE, align=PP_ALIGN.CENTER)


def slide_18_close(slide):
    """Close."""
    dark_bg_slide(slide)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK)
    add_rect(slide, 0, 0, 0.12, 7.5, fill_color=ACCENT_PURPLE)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=PRIMARY_BLUE)
    add_rect(slide, 8.5, 0, 5.0, 7.5, fill_color=RGBColor(20, 60, 140))

    # Main claim
    add_text_box(slide,
                 "The Most Composable ERP\nfor Africa",
                 0.6, 0.55, 9.5, 2.0,
                 font_size=40, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri")

    # Quote
    add_rect(slide, 0.6, 2.85, 8.8, 1.1,
             fill_color=RGBColor(20, 40, 90))
    add_rect(slide, 0.6, 2.85, 0.08, 1.1, fill_color=ACCENT_PURPLE)
    add_text_box(slide,
                 '"Compose Everything.\nDeploy Anywhere. Africa First."',
                 0.85, 2.9, 8.4, 0.95,
                 font_size=22, italic=True, bold=True,
                 color=LIGHT_BLUE, align=PP_ALIGN.LEFT, font_name="Calibri")

    # Stats row
    quick_stats = [
        ("155", "Modules"),
        ("588", "Cap. Models"),
        ("10",  "Connectors"),
        ("3",   "Tax APIs"),
    ]
    for i, (num, label) in enumerate(quick_stats):
        sx = 0.6 + i * 2.0
        add_rect(slide, sx, 4.2, 1.75, 1.3,
                 fill_color=RGBColor(26, 50, 120))
        add_text_box(slide, num, sx, 4.28, 1.75, 0.7,
                     font_size=36, bold=True, color=WARNING_ORANGE,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        add_text_box(slide, label, sx, 4.95, 1.75, 0.38,
                     font_size=12, color=LIGHT_BLUE,
                     align=PP_ALIGN.CENTER)

    # Contact / GitHub
    add_rect(slide, 0.6, 5.8, 8.5, 1.4,
             fill_color=RGBColor(15, 25, 55))
    add_text_box(slide, "Get in Touch",
                 0.75, 5.88, 3.0, 0.45,
                 font_size=14, bold=True, color=ACCENT_PURPLE)
    add_text_box(slide,
                 "GitHub:   github.com/your-org/pgappforge\n"
                 "Email:    founders@pgappforge.io\n"
                 "Website:  pgappforge.io",
                 0.75, 6.3, 8.0, 0.82,
                 font_size=13, color=LIGHT_BLUE,
                 font_name="Courier New")

    # Bottom bar
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_color=RGBColor(20, 60, 140))
    add_text_box(slide, "Confidential — For Qualified Investors Only",
                 0.5, 7.17, 12.0, 0.28,
                 font_size=10, color=RGBColor(148, 163, 184),
                 align=PP_ALIGN.CENTER)


# ─── Main ─────────────────────────────────────────────────────────────────────

def build_deck(output_path: str):
    prs = new_prs()

    builders = [
        slide_01_cover,
        slide_02_opportunity,
        slide_03_compliance,
        slide_04_incumbent_gap,
        slide_05_introducing,
        slide_06_composability_moat,
        slide_07_tech_moat,
        slide_08_sacco,
        slide_09_trade_finance,
        slide_10_library,
        slide_11_pdl_studio,
        slide_12_business_model,
        slide_13_gtm,
        slide_14_traction,
        slide_15_ask,
        slide_16_milestones,
        slide_17_why_now,
        slide_18_close,
    ]

    for fn in builders:
        s = blank_slide(prs)
        fn(s)
        print(f"  ✓  {fn.__name__}")

    prs.save(output_path)
    import os
    size = os.path.getsize(output_path)
    print(f"\nSaved: {output_path}")
    print(f"Size:  {size:,} bytes  ({size / 1024:.1f} KB)")
    return size


if __name__ == "__main__":
    out = "/Users/nyimbiodero/src/pjs/fab-ext/docs/composability/pgappforge-investor-pitch.pptx"
    build_deck(out)
