"""
Generate PgAppForge Composability Developer Conference Deck
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import os
import copy
from lxml import etree

# ─── Color palette ────────────────────────────────────────────────────────────
BLUE   = RGBColor(26,  86,  219)   # PgAppForge blue
DARK   = RGBColor(15,  23,  42)    # near-black
LIGHT  = RGBColor(248, 250, 252)   # off-white
PURPLE = RGBColor(124, 58,  237)   # accent / code
WHITE  = RGBColor(255, 255, 255)
GRAY   = RGBColor(100, 116, 139)
LGRAY  = RGBColor(226, 232, 240)
GREEN  = RGBColor(16,  185, 129)
CODE_BG = RGBColor(30,  30,  46)

# Slide dimensions: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    """Return the blank slide layout."""
    return prs.slide_layouts[6]


def add_slide(prs):
    return prs.slides.add_slide(blank_layout(prs))


def fill_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=18, bold=False, italic=False,
                color=DARK, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title(slide, title_text, top=Inches(0.4), color=WHITE, size=Pt(38), font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = font
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color
    return tb


def add_subtitle(slide, sub_text, top=Inches(1.5), color=LGRAY, size=Pt(22)):
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = sub_text
    run.font.name = "Calibri"
    run.font.size = size
    run.font.color.rgb = color
    return tb


def add_section_tag(slide, tag_text, tag_color=BLUE, left=Inches(0.6), top=Inches(0.15)):
    """Small colored pill label above the title."""
    rect = add_rect(slide, left, top, Inches(2.2), Inches(0.28), fill_color=tag_color)
    tf = rect.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = tag_text
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return rect


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=Pt(16), color=DARK, bullet_char="●  ",
                    line_spacing=1.3):
    """Add a list of bullet strings as separate paragraphs in one text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.name = "Calibri"
        run.font.size = font_size
        run.font.color.rgb = color
    return txBox


def add_code_box(slide, code_text, left, top, width, height, font_size=Pt(11)):
    """Dark rounded-corner code box with monospace white text."""
    # Background rounded rect
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    # Round corners via XML
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        else:
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 20000')

    # Text frame
    txBox = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.12),
        width - Inches(0.36), height - Inches(0.24)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    lines = code_text.strip('\n').split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = font_size
        run.font.color.rgb = WHITE
    return txBox


def add_divider_line(slide, top, color=BLUE):
    from pptx.util import Pt
    line = slide.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_slide_number(slide, num):
    add_textbox(slide, str(num),
                W - Inches(0.6), H - Inches(0.35),
                Inches(0.4), Inches(0.3),
                font_size=Pt(9), color=GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide builders ───────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide — dark blue background."""
    slide = add_slide(prs)
    fill_slide_bg(slide, DARK)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.07), fill_color=BLUE)

    # Decorative blue gradient block (right side)
    add_rect(slide, Inches(9.5), Inches(0), Inches(3.83), H, fill_color=RGBColor(20, 60, 160))

    # Small label
    add_section_tag(slide, "DEVELOPER CONFERENCE  2026", BLUE, Inches(0.6), Inches(1.2))

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.5), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "PgAppForge"
    run.font.name = "Calibri"
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Composability System"
    run2.font.name = "Calibri"
    run2.font.size = Pt(42)
    run2.font.bold = True
    run2.font.color.rgb = BLUE

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(8.5), Inches(0.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "11 Primitives That Let You Compose Any ERP Vertical"
    run3.font.name = "Calibri"
    run3.font.size = Pt(20)
    run3.font.color.rgb = LGRAY

    # Footer bar
    add_rect(slide, Inches(0), H - Inches(0.55), W, Inches(0.55), fill_color=RGBColor(10, 15, 30))
    add_textbox(slide, "Africa-First  ·  Python  ·  PostgreSQL  ·  Open Source",
                Inches(0.6), H - Inches(0.48), Inches(9), Inches(0.35),
                font_size=Pt(12), color=GRAY)


def slide_02_problem(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=RGBColor(220, 38, 38))
    add_section_tag(slide, "THE PROBLEM", RGBColor(220, 38, 38))
    add_title(slide, "ERP Systems Trap You", top=Inches(0.55), color=DARK, size=Pt(34))
    add_divider_line(slide, Inches(1.55), color=RGBColor(220, 38, 38))

    items = [
        "Customization hell: touching core code = upgrade nightmares",
        "Plugin lock-in: Plugin A can't extend Plugin B's models without forking",
        "Cross-domain gaps: loyalty + mobile money + tax need custom glue",
    ]
    add_bullet_list(slide, items, Inches(0.8), Inches(1.75), Inches(11.5), Inches(4),
                    font_size=Pt(20), color=DARK, bullet_char="✗  ")
    add_slide_number(slide, 2)


def slide_03_primitives(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "OVERVIEW", BLUE)
    add_title(slide, "11 Composition Primitives", top=Inches(0.55), color=DARK, size=Pt(34))
    add_divider_line(slide, Inches(1.55))

    left_items = [
        "Event Router (glob patterns)",
        "Model Mixin Registry",
        "Sub-workflow Composition",
        "Rule → Event Bridge",
        "Permission Algebra",
        "PDL Schema Extension",
    ]
    right_items = [
        "Semantic Metric Registry",
        "View Slot Injection",
        "AI Composable Pipeline",
        "Cross-Tenant Aggregation",
        "GraphQL Federation",
    ]

    # Left column header
    add_textbox(slide, "Core Primitives", Inches(0.8), Inches(1.75), Inches(5.5), Inches(0.45),
                font_size=Pt(13), bold=True, color=BLUE)
    add_bullet_list(slide, left_items, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.0),
                    font_size=Pt(16), color=DARK)

    # Right column header
    add_textbox(slide, "Advanced Primitives", Inches(7.0), Inches(1.75), Inches(5.5), Inches(0.45),
                font_size=Pt(13), bold=True, color=PURPLE)
    add_bullet_list(slide, right_items, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.0),
                    font_size=Pt(16), color=DARK)

    # Vertical divider
    add_rect(slide, Inches(6.65), Inches(1.75), Pt(1.5), Inches(5.0), fill_color=LGRAY)
    add_slide_number(slide, 3)


def slide_04_event_router(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "PRIMITIVE 1 — EVENT ROUTER", BLUE)
    add_title(slide, "Event Router: React to Anything", top=Inches(0.55), color=DARK, size=Pt(32))
    add_divider_line(slide, Inches(1.55))

    code = """\
@on_event('finance.ar.invoice.*')
def sync_to_tax(event_type, payload, tenant_id):
    etims.submit(payload['invoice_id'])

emit('finance.ar.invoice.approved',
     {'invoice_id': inv.id}, tenant_id=tid, session=session)"""
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(2.5))

    tags = ["Glob patterns", "Exception isolation", "Durable", "Multi-worker"]
    for i, tag in enumerate(tags):
        x = Inches(0.7) + i * Inches(3.0)
        pill = add_rect(slide, x, Inches(4.55), Inches(2.6), Inches(0.4), fill_color=RGBColor(219, 234, 254))
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "✓  " + tag
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = BLUE

    add_textbox(slide,
                "Events are durable, namespaced by plugin+domain, and processed in isolation per handler.",
                Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.6),
                font_size=Pt(13), color=GRAY)
    add_slide_number(slide, 4)


def slide_05_mixin_registry(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "PRIMITIVE 2 — MODEL MIXIN REGISTRY", BLUE)
    add_title(slide, "Model Mixin Registry: Extend Any Model", top=Inches(0.55), color=DARK, size=Pt(32))
    add_divider_line(slide, Inches(1.55))

    code = """\
class LCMixin:
    letter_of_credit_id = sa.Column(sa.String(36))
    lc_expiry_date = sa.Column(sa.Date)

register_mixin(
    'pgappforge.plugins.erp.finance.ar.models.ARInvoice',
    LCMixin, priority=10
)"""
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(2.8))

    add_rect(slide, Inches(0.7), Inches(4.75), Inches(11.9), Inches(0.65),
             fill_color=RGBColor(240, 253, 244))
    add_textbox(slide,
                "ℹ  Applied before SQLAlchemy mapper configuration — no monkey-patching, no model forking.",
                Inches(0.9), Inches(4.82), Inches(11.5), Inches(0.5),
                font_size=Pt(13), color=RGBColor(21, 128, 61))
    add_slide_number(slide, 5)


def slide_06_permission_algebra(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=PURPLE)
    add_section_tag(slide, "PRIMITIVE 5 — PERMISSION ALGEBRA", PURPLE)
    add_title(slide, "Permission Algebra: Business Logic, Not Code", top=Inches(0.55), color=DARK, size=Pt(30))
    add_divider_line(slide, Inches(1.55), color=PURPLE)

    code = """\
approve_loan = AllOf(
    HasRole('loan_officer'),
    AnyOf(
        HasRole('credit_committee'),
        HasPermission('credit.emergency_override'),
    ),
)

@require_policy(approve_loan)
def approve(self, loan_id): ..."""
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(3.2))

    add_textbox(slide,
                "Composable policy objects — no permission spreadsheet, no spaghetti if-else chains.",
                Inches(0.7), Inches(5.2), Inches(11.9), Inches(0.6),
                font_size=Pt(14), color=GRAY)
    add_slide_number(slide, 6)


def slide_07_ai_pipeline(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=PURPLE)
    add_section_tag(slide, "PRIMITIVE 9 — AI COMPOSABLE PIPELINE", PURPLE)
    add_title(slide, "AI Composable Pipeline: The Runnable Protocol", top=Inches(0.55), color=DARK, size=Pt(29))
    add_divider_line(slide, Inches(1.55), color=PURPLE)

    code = """\
pipeline = (
    SQLStep("SELECT * FROM fin_ar_invoice WHERE status='OVERDUE'")
    | LLMStep(system="Summarise for the CFO.")
    | FormatStep("CFO Report:\\n{value}")
)
result = pipeline.invoke({"tenant_id": tid}, session=session)"""
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(2.6))

    add_rect(slide, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.55),
             fill_color=RGBColor(245, 243, 255))
    add_textbox(slide,
                "Inspired by LangChain LCEL  ·  Works with any LiteLLM model  ·  Tenant-scoped execution",
                Inches(0.9), Inches(4.65), Inches(11.5), Inches(0.45),
                font_size=Pt(13), color=PURPLE)
    add_slide_number(slide, 7)


def slide_08_pdl(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "PRIMITIVE 6 — PDL SCHEMA EXTENSION", BLUE)
    add_title(slide, "PDL Schema Extension", top=Inches(0.55), color=DARK, size=Pt(34))
    add_divider_line(slide, Inches(1.55))

    code = """\
entities:
  - name: LoyalCustomer
    table: crm_loyal_customer
    extends: pgappforge.plugins.erp.crm.prm.models.PartnerAccount
    fields:
      - name: points_balance
        type: integer
      - name: tier
        type: string
        default: "bronze" """
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(3.2))

    add_textbox(slide,
                "One YAML file drives: model class · Alembic migration · REST API · OpenAPI schema · tests.",
                Inches(0.7), Inches(5.2), Inches(11.9), Inches(0.6),
                font_size=Pt(14), color=GRAY)
    add_slide_number(slide, 8)


def slide_09_metrics(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "PRIMITIVE 7 — SEMANTIC METRIC REGISTRY", BLUE)
    add_title(slide, "Semantic Metric Registry", top=Inches(0.55), color=DARK, size=Pt(34))
    add_divider_line(slide, Inches(1.55))

    code = """\
register_metric(Metric(
    name='finance.ar.revenue',
    model_path='...ARInvoice',
    field='total_amount_cents',
    agg='sum',  # additive: safe to sum across groups
))

results = query_metrics(
    ['finance.ar.revenue', 'crm.deals_won', 'hcm.headcount_cost'],
    group_by=['tenant_id'], session=session,
)"""
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(3.5))
    add_slide_number(slide, 9)


def slide_10_view_slots(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "PRIMITIVE 8 — VIEW SLOT INJECTION", BLUE)
    add_title(slide, "View Slot Injection", top=Inches(0.55), color=DARK, size=Pt(34))
    add_divider_line(slide, Inches(1.55))

    code = """\
@slot_provider('customer.detail.sidebar', priority=10)
def loyalty_widget(context: dict) -> str:
    return f'<div>Points: {get_points(context["customer_id"])}</div>'

# In your Jinja2 template:
# {{ render_slot('customer.detail.sidebar', ctx) }}"""
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(2.7))

    add_textbox(slide,
                "Any plugin can inject UI widgets into any named slot — zero coupling to the host view.",
                Inches(0.7), Inches(4.7), Inches(11.9), Inches(0.6),
                font_size=Pt(14), color=GRAY)
    add_slide_number(slide, 10)


def slide_11_sacco(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=GREEN)
    add_section_tag(slide, "REAL-WORLD EXAMPLE — SACCO PLATFORM", GREEN)
    add_title(slide, "Real-World Composition: SACCO Platform", top=Inches(0.55), color=DARK, size=Pt(30))
    add_divider_line(slide, Inches(1.55), color=GREEN)

    steps = [
        "Member joins → KYC workflow  (sub-workflow composition)",
        "Loan application → AllOf(loan_officer, AnyOf(manager, committee))  (permission algebra)",
        "Loan approved → emit('sacco.loan.approved')  (event router)",
        "EventRouter → eTIMS + MTN MoMo + SMS  (3 handlers, one event)",
        "Repayment → Loyalty points → Dashboard slot updates  (metric + view slots)",
    ]
    for i, step in enumerate(steps):
        y = Inches(1.8) + i * Inches(0.88)
        num_rect = add_rect(slide, Inches(0.7), y, Inches(0.45), Inches(0.45),
                            fill_color=BLUE)
        tf = num_rect.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = WHITE

        add_textbox(slide, step, Inches(1.3), y + Inches(0.02), Inches(11.0), Inches(0.45),
                    font_size=Pt(15), color=DARK)
    add_slide_number(slide, 11)


def slide_12_588(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "VISUAL DESIGNER", BLUE)
    add_title(slide, "588 Capability Models — Visual Designer", top=Inches(0.55), color=DARK, size=Pt(30))
    add_divider_line(slide, Inches(1.55))

    items = [
        "Import any of 588 pre-built capability models from 8 domains",
        "Draw FK relationships (auto-layout with barycenter crossing minimization)",
        "Generate: model + migration + REST API + tests + Dockerfile + K8s",
        "One YAML schema → complete deployable service",
    ]
    add_bullet_list(slide, items, Inches(0.8), Inches(1.85), Inches(11.5), Inches(4.5),
                    font_size=Pt(19), color=DARK, bullet_char="→  ")

    # Bottom bar
    add_rect(slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.6),
             fill_color=RGBColor(219, 234, 254))
    add_textbox(slide,
                "flask forge designer --open    |    flask forge gen pdl schema.pdl.yaml --with-docker --with-k8s",
                Inches(0.9), Inches(6.47), Inches(11.5), Inches(0.45),
                font_name="Courier New", font_size=Pt(12), color=BLUE)
    add_slide_number(slide, 12)


def slide_13_comparison(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "COMPETITIVE ANALYSIS", BLUE)
    add_title(slide, "Comparison vs Best-in-Class", top=Inches(0.55), color=DARK, size=Pt(34))
    add_divider_line(slide, Inches(1.55))

    headers = ["Feature", "PgAppForge", "Odoo", "JHipster", "SAP B1"]
    rows = [
        ["Model inheritance",  "✅ Mixin registry",   "✅ _inherit",  "❌", "❌"],
        ["Policy algebra",     "✅ AllOf/AnyOf",       "❌ flat ACL",  "❌", "❌"],
        ["AI pipeline",        "✅ Runnable",          "❌",           "❌", "❌"],
        ["Africa connectors",  "✅ 10 connectors",     "❌",           "❌", "❌"],
        ["GraphQL federation", "✅ Apollo v2",         "❌",           "❌", "❌"],
    ]

    col_widths = [Inches(2.8), Inches(2.6), Inches(1.9), Inches(1.9), Inches(1.9)]
    col_x = [Inches(0.7)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    row_h = Inches(0.55)
    table_top = Inches(1.75)

    # Header row
    for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        hdr_rect = add_rect(slide, cx, table_top, cw, row_h,
                            fill_color=BLUE if ci == 0 else RGBColor(30, 58, 138))
        tf = hdr_rect.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = WHITE

    for ri, row in enumerate(rows):
        row_top = table_top + row_h + ri * row_h
        for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_widths)):
            bg = RGBColor(239, 246, 255) if ri % 2 == 0 else WHITE
            if ci == 0:
                bg = RGBColor(219, 234, 254)
            cell_rect = add_rect(slide, cx, row_top, cw, row_h, fill_color=bg,
                                 line_color=LGRAY, line_width=Pt(0.5))
            tf = cell_rect.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = cell
            run.font.name = "Calibri"
            run.font.size = Pt(12)
            run.font.color.rgb = DARK if ci > 0 else RGBColor(30, 58, 138)
            run.font.bold = (ci == 0)

    add_slide_number(slide, 13)


def slide_14_graphql(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, LIGHT)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=PURPLE)
    add_section_tag(slide, "PRIMITIVES 10 & 11 — GRAPHQL + CROSS-TENANT", PURPLE)
    add_title(slide, "GraphQL Federation + Cross-Tenant", top=Inches(0.55), color=DARK, size=Pt(31))
    add_divider_line(slide, Inches(1.55), color=PURPLE)

    # Left: federation
    add_textbox(slide, "GraphQL Federation", Inches(0.7), Inches(1.75), Inches(5.8), Inches(0.45),
                font_size=Pt(14), bold=True, color=PURPLE)
    code_fed = """\
@federated_type(key='id', plugin='finance.ar')
class ARInvoice:
    id: str
    total_amount_cents: int"""
    add_code_box(slide, code_fed, Inches(0.7), Inches(2.25), Inches(5.8), Inches(1.6), font_size=Pt(11))

    # Right: cross-tenant
    add_textbox(slide, "Cross-Tenant Aggregation", Inches(7.0), Inches(1.75), Inches(5.8), Inches(0.45),
                font_size=Pt(14), bold=True, color=BLUE)
    code_ct = """\
with SystemSession(session,
        caller_user_id='admin',
        reason='billing'):
    summary = aggregator.get_platform_summary(
        session)"""
    add_code_box(slide, code_ct, Inches(7.0), Inches(2.25), Inches(5.8), Inches(1.6), font_size=Pt(11))

    add_textbox(slide,
                "Federation: expose any plugin type to a unified GraphQL schema.\n"
                "Cross-tenant: platform-level aggregation with full audit trail.",
                Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.9),
                font_size=Pt(14), color=GRAY)
    add_slide_number(slide, 14)


def slide_15_getting_started(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill_color=BLUE)
    add_section_tag(slide, "GET STARTED", BLUE)
    add_title(slide, "Getting Started", top=Inches(0.55), color=WHITE, size=Pt(36))
    add_divider_line(slide, Inches(1.55))

    code = """\
pip install pgappforge

flask forge gen pdl schema.pdl.yaml --with-docker --with-k8s

flask forge designer --open"""
    add_code_box(slide, code, Inches(0.7), Inches(1.75), Inches(11.9), Inches(1.8))

    steps = [
        "Define your domain schema in PDL YAML",
        "Generate models, migrations, REST API, and tests in one command",
        "Open the Visual Designer — drag, drop, connect, deploy",
        "Compose plugins with event routing, mixins, and policy algebra",
    ]
    add_bullet_list(slide, steps, Inches(0.8), Inches(3.75), Inches(11.5), Inches(2.5),
                    font_size=Pt(16), color=LGRAY, bullet_char="→  ")

    # Links row
    add_rect(slide, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.55),
             fill_color=RGBColor(20, 40, 80))
    add_textbox(slide,
                "docs.pgappforge.io    |    github.com/pgappforge",
                Inches(0.9), Inches(6.52), Inches(11.5), Inches(0.45),
                font_size=Pt(14), color=BLUE, align=PP_ALIGN.CENTER)
    add_slide_number(slide, 15)


def slide_16_thankyou(prs):
    slide = add_slide(prs)
    fill_slide_bg(slide, DARK)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.07), fill_color=BLUE)

    # Decorative side block
    add_rect(slide, Inches(9.5), Inches(0), Inches(3.83), H, fill_color=RGBColor(20, 60, 160))

    add_textbox(slide, "Thank You",
                Inches(0.6), Inches(2.0), Inches(8.5), Inches(1.4),
                font_name="Calibri", font_size=Pt(58), bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    add_textbox(slide, "Compose Everything. Deploy Anywhere. Africa First.",
                Inches(0.6), Inches(3.55), Inches(8.5), Inches(0.8),
                font_name="Calibri", font_size=Pt(22), color=BLUE)

    add_textbox(slide, "Q & A",
                Inches(0.6), Inches(4.5), Inches(3), Inches(0.8),
                font_name="Calibri", font_size=Pt(30), bold=True, color=LGRAY)

    add_rect(slide, Inches(0), H - Inches(0.55), W, Inches(0.55), fill_color=RGBColor(10, 15, 30))
    add_textbox(slide,
                "Africa-First  ·  Python  ·  PostgreSQL  ·  Open Source  ·  pgappforge.io",
                Inches(0.6), H - Inches(0.48), Inches(9), Inches(0.35),
                font_size=Pt(11), color=GRAY)


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_primitives(prs)
    slide_04_event_router(prs)
    slide_05_mixin_registry(prs)
    slide_06_permission_algebra(prs)
    slide_07_ai_pipeline(prs)
    slide_08_pdl(prs)
    slide_09_metrics(prs)
    slide_10_view_slots(prs)
    slide_11_sacco(prs)
    slide_12_588(prs)
    slide_13_comparison(prs)
    slide_14_graphql(prs)
    slide_15_getting_started(prs)
    slide_16_thankyou(prs)

    output_path = '/Users/nyimbiodero/src/pjs/fab-ext/docs/composability/pgappforge-composability-developer.pptx'
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    size = os.path.getsize(output_path)
    print(f'Saved to {output_path}')
    print(f'developer deck: OK, size: {size} bytes')
    return size


if __name__ == '__main__':
    main()
